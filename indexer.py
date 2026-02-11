"""
Email indexer for Outlook emails.
Fetches emails from Microsoft 365 via Graph API, generates embeddings,
and stores them in ChromaDB for semantic search.
"""

import json
import logging
import re
import time

# Suppress model loading warnings
import os
os.environ["TRANSFORMERS_VERBOSITY"] = "error"
os.environ["HF_HUB_DISABLE_PROGRESS_BARS"] = "1"
logging.getLogger("sentence_transformers").setLevel(logging.ERROR)
logging.getLogger("transformers").setLevel(logging.ERROR)

import warnings
warnings.filterwarnings("ignore")
import base64
import io
from datetime import datetime, timezone
from typing import Optional, List, Dict, Any, Generator
from bs4 import BeautifulSoup
import msal
import requests
import chromadb
from chromadb.config import Settings as ChromaSettings
from sentence_transformers import SentenceTransformer

# Attachment parsing imports (optional - gracefully handle if not installed)
try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import load_workbook
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from config import (
    get_settings,
    IndexState,
    DATA_DIR,
    Settings,
)


# Microsoft Graph API endpoints
GRAPH_BASE_URL = "https://graph.microsoft.com/v1.0"
SCOPES = ["https://graph.microsoft.com/Mail.Read"]

# Supported attachment extensions
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.csv', '.json', '.md', '.html', '.htm'}


def extract_text_from_attachment(content_bytes: bytes, filename: str) -> str:
    """Extract text content from an attachment based on its file type."""
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    ext = '.' + ext

    try:
        if ext == '.pdf' and HAS_PDF:
            return _extract_pdf_text(content_bytes)
        elif ext == '.docx' and HAS_DOCX:
            return _extract_docx_text(content_bytes)
        elif ext == '.xlsx' and HAS_XLSX:
            return _extract_xlsx_text(content_bytes)
        elif ext == '.pptx' and HAS_PPTX:
            return _extract_pptx_text(content_bytes)
        elif ext in {'.txt', '.csv', '.json', '.md'}:
            return _extract_plain_text(content_bytes)
        elif ext in {'.html', '.htm'}:
            return _extract_html_text(content_bytes)
    except Exception as e:
        print(f"    Warning: Could not extract text from {filename}: {e}")
    return ""


def _extract_pdf_text(content_bytes: bytes) -> str:
    """Extract text from PDF."""
    reader = PdfReader(io.BytesIO(content_bytes))
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return "\n".join(text_parts)


def _extract_docx_text(content_bytes: bytes) -> str:
    """Extract text from Word document."""
    doc = DocxDocument(io.BytesIO(content_bytes))
    return "\n".join(para.text for para in doc.paragraphs if para.text)


def _extract_xlsx_text(content_bytes: bytes) -> str:
    """Extract text from Excel spreadsheet."""
    wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
    text_parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                text_parts.append(row_text)
    return "\n".join(text_parts)


def _extract_pptx_text(content_bytes: bytes) -> str:
    """Extract text from PowerPoint presentation."""
    prs = Presentation(io.BytesIO(content_bytes))
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
    return "\n".join(text_parts)


def _extract_plain_text(content_bytes: bytes) -> str:
    """Extract plain text, trying different encodings."""
    for encoding in ['utf-8', 'latin-1', 'cp1252']:
        try:
            return content_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _extract_html_text(content_bytes: bytes) -> str:
    """Extract text from HTML."""
    html = _extract_plain_text(content_bytes)
    if html:
        soup = BeautifulSoup(html, "html.parser")
        for element in soup(["script", "style", "head", "meta"]):
            element.decompose()
        return soup.get_text(separator=" ")
    return ""


class MSGraphAuth:
    """Handles Microsoft Graph API authentication using device code flow."""

    def __init__(self, client_id: str, tenant_id: str = "common"):
        self.client_id = client_id
        self.tenant_id = tenant_id
        self.authority = f"https://login.microsoftonline.com/{tenant_id}"
        self.app = msal.PublicClientApplication(
            client_id,
            authority=self.authority,
        )
        self._access_token: Optional[str] = None
        self._token_expiry: Optional[datetime] = None

    def get_access_token(self) -> str:
        """Get a valid access token, refreshing if needed."""
        # Check if we have a valid cached token
        if self._access_token and self._token_expiry:
            if datetime.now(timezone.utc) < self._token_expiry:
                return self._access_token

        # Try to get token from cache
        accounts = self.app.get_accounts()
        if accounts:
            result = self.app.acquire_token_silent(SCOPES, account=accounts[0])
            if result and "access_token" in result:
                self._cache_token(result)
                return self._access_token

        # No cached token, need to authenticate
        flow = self.app.initiate_device_flow(scopes=SCOPES)
        if "user_code" not in flow:
            raise Exception(f"Failed to create device flow: {flow.get('error_description', 'Unknown error')}")

        print("\n" + "=" * 60)
        print("Microsoft Graph API Authentication Required")
        print("=" * 60)
        print(f"\nTo sign in, visit: {flow['verification_uri']}")
        print(f"Enter code: {flow['user_code']}")
        print("\nWaiting for authentication...")

        result = self.app.acquire_token_by_device_flow(flow)

        if "access_token" not in result:
            raise Exception(f"Authentication failed: {result.get('error_description', 'Unknown error')}")

        self._cache_token(result)
        print("Authentication successful!\n")
        return self._access_token

    def _cache_token(self, result: dict) -> None:
        """Cache the access token and expiry."""
        self._access_token = result["access_token"]
        expires_in = result.get("expires_in", 3600)
        self._token_expiry = datetime.now(timezone.utc).replace(
            microsecond=0
        ) + __import__("datetime").timedelta(seconds=expires_in - 60)


class EmailFetcher:
    """Fetches emails from Microsoft Graph API."""

    def __init__(self, auth: MSGraphAuth):
        self.auth = auth

    def _make_request(self, url: str, params: Optional[dict] = None) -> dict:
        """Make an authenticated request to Graph API."""
        headers = {
            "Authorization": f"Bearer {self.auth.get_access_token()}",
            "Content-Type": "application/json",
        }
        response = requests.get(url, headers=headers, params=params)
        response.raise_for_status()
        return response.json()

    def get_mail_folders(self) -> List[Dict[str, Any]]:
        """Get all mail folders recursively."""
        folders = []

        def fetch_folders(parent_id: Optional[str] = None):
            if parent_id:
                url = f"{GRAPH_BASE_URL}/me/mailFolders/{parent_id}/childFolders"
            else:
                url = f"{GRAPH_BASE_URL}/me/mailFolders"

            while url:
                data = self._make_request(url)
                for folder in data.get("value", []):
                    folders.append({
                        "id": folder["id"],
                        "name": folder["displayName"],
                        "total_count": folder.get("totalItemCount", 0),
                    })
                    # Recursively get child folders
                    if folder.get("childFolderCount", 0) > 0:
                        fetch_folders(folder["id"])

                url = data.get("@odata.nextLink")

        fetch_folders()
        return folders

    def get_emails(
        self,
        folder_id: str,
        since: Optional[str] = None,
        batch_size: int = 50,
        include_attachments: bool = True,
    ) -> Generator[Dict[str, Any], None, None]:
        """
        Fetch emails from a folder with pagination.

        Args:
            folder_id: The folder ID to fetch from
            since: ISO timestamp to fetch emails received after (for incremental indexing)
            batch_size: Number of emails per API call (max 1000)
            include_attachments: Whether to fetch and extract text from attachments

        Yields:
            Email dictionaries with extracted content
        """
        url = f"{GRAPH_BASE_URL}/me/mailFolders/{folder_id}/messages"

        params = {
            "$top": min(batch_size, 1000),
            "$select": "id,subject,body,from,toRecipients,ccRecipients,receivedDateTime,conversationId,hasAttachments",
            "$orderby": "receivedDateTime desc",
        }

        if since:
            params["$filter"] = f"receivedDateTime gt {since}"

        while url:
            data = self._make_request(url, params if "?" not in url else None)

            for email in data.get("value", []):
                try:
                    email_data = self._extract_email_data(email)

                    # Fetch attachments if present and requested
                    if include_attachments and email.get("hasAttachments"):
                        attachment_text = self._get_attachment_text(email["id"])
                        if attachment_text:
                            email_data["body"] += "\n\n[ATTACHMENTS]\n" + attachment_text

                    yield email_data
                except Exception as e:
                    subject = email.get("subject", "unknown")[:50]
                    print(f"    Warning: Skipping email '{subject}': {e}")
                    continue

            url = data.get("@odata.nextLink")
            params = None  # nextLink includes parameters

    def _get_attachment_text(self, email_id: str) -> str:
        """Fetch and extract text from email attachments."""
        url = f"{GRAPH_BASE_URL}/me/messages/{email_id}/attachments"

        try:
            data = self._make_request(url)
        except Exception as e:
            print(f"    Warning: Could not fetch attachments: {e}")
            return ""

        text_parts = []

        for attachment in data.get("value", []):
            # Only process file attachments (not item attachments like nested emails)
            if attachment.get("@odata.type") != "#microsoft.graph.fileAttachment":
                continue

            filename = attachment.get("name", "")
            ext = '.' + filename.lower().split('.')[-1] if '.' in filename else ''

            # Skip unsupported file types
            if ext not in SUPPORTED_EXTENSIONS:
                continue

            # Get content (base64 encoded)
            content_b64 = attachment.get("contentBytes", "")
            if not content_b64:
                continue

            try:
                content_bytes = base64.b64decode(content_b64)
                text = extract_text_from_attachment(content_bytes, filename)
                if text:
                    text_parts.append(f"[{filename}]\n{text}")
            except Exception as e:
                print(f"    Warning: Could not process {filename}: {e}")

        return "\n\n".join(text_parts)

    def _extract_email_data(self, email: dict) -> Dict[str, Any]:
        """Extract and clean email data."""
        # Extract body text, stripping HTML (handle None values)
        body = email.get("body") or {}
        body_content = body.get("content", "") or ""
        body_type = body.get("contentType", "text") or "text"

        if body_type.lower() == "html":
            body_text = self._strip_html(body_content)
        else:
            body_text = body_content

        # Clean up whitespace
        body_text = re.sub(r'\s+', ' ', body_text).strip()

        # Extract sender (handle None values)
        from_field = email.get("from") or {}
        from_data = from_field.get("emailAddress") or {}
        sender = from_data.get("address", "") or ""
        sender_name = from_data.get("name", "") or ""

        # Extract recipients (handle None values)
        to_recipients = [
            (r.get("emailAddress") or {}).get("address", "") or ""
            for r in (email.get("toRecipients") or [])
        ]
        cc_recipients = [
            (r.get("emailAddress") or {}).get("address", "") or ""
            for r in (email.get("ccRecipients") or [])
        ]

        return {
            "id": email.get("id", ""),
            "subject": email.get("subject", ""),
            "body": body_text,
            "sender": sender,
            "sender_name": sender_name,
            "to": to_recipients,
            "cc": cc_recipients,
            "received_datetime": email.get("receivedDateTime", ""),
            "conversation_id": email.get("conversationId", ""),
        }

    def _strip_html(self, html: str) -> str:
        """Strip HTML tags and extract text content."""
        soup = BeautifulSoup(html, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style", "head", "meta"]):
            element.decompose()

        # Get text
        text = soup.get_text(separator=" ")
        return text


class EmailIndexer:
    """Indexes emails into ChromaDB with embeddings."""

    def __init__(self, settings: Optional[Settings] = None):
        self.settings = settings or get_settings()

        # Initialize embedding model
        print(f"Loading embedding model: {self.settings.embedding_model}")
        self.embedder = SentenceTransformer(self.settings.embedding_model)

        # Initialize ChromaDB
        self.chroma_client = chromadb.PersistentClient(
            path=str(DATA_DIR / "chromadb"),
            settings=ChromaSettings(anonymized_telemetry=False),
        )
        self.collection = self.chroma_client.get_or_create_collection(
            name=self.settings.chroma_collection_name,
            metadata={"hnsw:space": "cosine"},
        )

    def index_emails(
        self,
        emails: List[Dict[str, Any]],
        folder_name: str,
        batch_size: int = 100,
    ) -> int:
        """
        Index a list of emails into ChromaDB.

        Args:
            emails: List of email dictionaries
            folder_name: Name of the source folder
            batch_size: Number of emails to embed at once

        Returns:
            Number of emails indexed
        """
        if not emails:
            return 0

        indexed = 0

        for i in range(0, len(emails), batch_size):
            batch = emails[i:i + batch_size]

            # Create text for embedding (subject + body)
            texts = []
            valid_batch = []
            for email in batch:
                try:
                    subject = email.get('subject') or ''
                    body = email.get('body') or ''
                    text = f"Subject: {subject}\n\n{body}"
                    # Truncate very long emails (sentence-transformers has token limits)
                    if len(text) > 8000:
                        text = text[:8000] + "..."
                    texts.append(text)
                    valid_batch.append(email)
                except Exception as e:
                    print(f"    Warning: Skipping malformed email: {e}")
                    continue

            if not valid_batch:
                continue

            batch = valid_batch

            # Generate embeddings
            embeddings = self.embedder.encode(texts, show_progress_bar=False)

            # Prepare data for ChromaDB
            ids = [email.get("id") or f"unknown_{i}" for i, email in enumerate(batch)]
            metadatas = [
                {
                    "subject": (email.get("subject") or "")[:500],  # Truncate for metadata
                    "sender": email.get("sender") or "",
                    "sender_name": email.get("sender_name") or "",
                    "to": json.dumps((email.get("to") or [])[:10]),  # Limit recipients
                    "received_datetime": email.get("received_datetime") or "",
                    "folder": folder_name,
                    "conversation_id": email.get("conversation_id") or "",
                }
                for email in batch
            ]
            documents = texts

            # Upsert to collection (handles duplicates)
            self.collection.upsert(
                ids=ids,
                embeddings=embeddings.tolist(),
                metadatas=metadatas,
                documents=documents,
            )

            indexed += len(batch)
            print(f"  Indexed {indexed} emails...")

        return indexed

    def get_stats(self) -> Dict[str, Any]:
        """Get indexing statistics."""
        return {
            "total_emails": self.collection.count(),
            "collection_name": self.settings.chroma_collection_name,
        }


def run_indexer(
    settings: Optional[Settings] = None,
    full_reindex: bool = False,
    verbose: bool = True,
) -> Dict[str, Any]:
    """
    Main indexing function. Fetches emails from MS Graph and indexes them.

    Args:
        settings: Application settings (uses global settings if not provided)
        full_reindex: If True, ignore incremental state and reindex everything
        verbose: Print progress messages

    Returns:
        Dictionary with indexing results
    """
    settings = settings or get_settings()

    if not settings.ms_graph_client_id:
        raise ValueError(
            "MS Graph client ID not configured. "
            "Please set ms_graph_client_id in settings.json"
        )

    # Initialize components
    auth = MSGraphAuth(settings.ms_graph_client_id, settings.ms_graph_tenant_id)
    fetcher = EmailFetcher(auth)
    indexer = EmailIndexer(settings)
    state = IndexState.load()

    results = {
        "folders_processed": 0,
        "emails_indexed": 0,
        "errors": [],
    }

    # Get incremental timestamp
    since = None if full_reindex else state.last_indexed.get("msgraph")

    if verbose:
        if since:
            print(f"Incremental index from: {since}")
        else:
            print("Full index (no previous state)")

    # Get all folders
    if verbose:
        print("Fetching mail folders...")
    folders = fetcher.get_mail_folders()
    if verbose:
        print(f"Found {len(folders)} folders")

    # Track latest email timestamp for incremental updates
    latest_timestamp = since

    # Index each folder
    for folder in folders:
        if verbose:
            print(f"\nProcessing folder: {folder['name']} ({folder['total_count']} items)")

        try:
            emails = list(fetcher.get_emails(folder["id"], since=since))

            if emails:
                try:
                    indexed = indexer.index_emails(emails, folder["name"])
                    results["emails_indexed"] += indexed
                except Exception as e:
                    error_msg = f"Error indexing folder {folder['name']}: {e}"
                    results["errors"].append(error_msg)
                    if verbose:
                        print(f"  ERROR during indexing: {e}")
                    # Continue to next folder

                # Track latest timestamp
                for email in emails:
                    ts = email.get("received_datetime", "") or ""
                    if ts and (not latest_timestamp or ts > latest_timestamp):
                        latest_timestamp = ts

            results["folders_processed"] += 1

        except Exception as e:
            error_msg = f"Error processing folder {folder['name']}: {e}"
            results["errors"].append(error_msg)
            if verbose:
                print(f"  ERROR: {e}")

    # Update state
    if latest_timestamp:
        state.update_source(
            "msgraph",
            latest_timestamp,
            indexer.get_stats()["total_emails"],
        )

    if verbose:
        print(f"\n{'=' * 60}")
        print(f"Indexing complete!")
        print(f"Folders processed: {results['folders_processed']}")
        print(f"Emails indexed: {results['emails_indexed']}")
        print(f"Total emails in database: {indexer.get_stats()['total_emails']}")
        if results["errors"]:
            print(f"Errors: {len(results['errors'])}")

    return results


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Index Outlook emails for semantic search")
    parser.add_argument("--full", action="store_true", help="Full reindex (ignore incremental state)")
    parser.add_argument("--quiet", action="store_true", help="Suppress progress messages")
    args = parser.parse_args()

    try:
        run_indexer(full_reindex=args.full, verbose=not args.quiet)
    except KeyboardInterrupt:
        print("\nIndexing interrupted.")
    except Exception as e:
        print(f"Error: {e}")
        raise
